"""
BPL Presentation Generator — Professional Grade
Why Bangladesh Premier League Fails to Attract Urban Youth
Created for: Textile Department, Batch 232
© Copyright by Tahorim Somrat
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
GREEN       = RGBColor(0x00, 0xC8, 0x53)
DARK_GREEN  = RGBColor(0x00, 0x66, 0x2A)
RED         = RGBColor(0xFF, 0x17, 0x44)
GOLD        = RGBColor(0xFF, 0xD7, 0x00)
ORANGE      = RGBColor(0xFF, 0x77, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x07, 0x0B, 0x16)
DARK_CARD   = RGBColor(0x0F, 0x16, 0x2A)
CARD2       = RGBColor(0x12, 0x1C, 0x38)
ACCENT_BLUE = RGBColor(0x00, 0x8B, 0xFF)
TEAL        = RGBColor(0x00, 0xD4, 0xAA)
LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE8)
MID_GRAY    = RGBColor(0x55, 0x66, 0x88)

OUT_PATH  = r"C:\Users\mdsom\Desktop\Tahorim AI\Abid\BPL_Presentation_v2.pptx"
ABID_DIR  = r"C:\Users\mdsom\Desktop\Tahorim AI\Abid"

IMG_SLIDE1   = os.path.join(ABID_DIR, "bpl_slide1_bg.png")
IMG_INTL     = os.path.join(ABID_DIR, "international_stars.png")
IMG_BD       = os.path.join(ABID_DIR, "bangladesh_players.png")
IMG_STADIUM  = os.path.join(ABID_DIR, "bangladesh_stadium.png")
IMG_USERPHOTO= os.path.join(ABID_DIR, "user_split_image.jpg")

# ─── Core Helpers ────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name="Segoe UI"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(font_size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    run.font.name       = font_name
    return txb

def add_multiline(slide, lines, l, t, w, h, default_size=16, default_bold=False,
                  default_color=WHITE, align=PP_ALIGN.LEFT, spacing_after=8):
    """lines = list of (text, font_size, bold, color) tuples"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (txt, sz, bold, clr) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text           = txt
        run.font.size      = Pt(sz if sz else default_size)
        run.font.bold      = bold if bold is not None else default_bold
        run.font.color.rgb = clr if clr else default_color
        run.font.name      = "Segoe UI"
    return txb

def add_gradient_rect(slide, l, t, w, h, c1, c2, angle=90):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    sp   = shape._element
    spPr = sp.find(qn('p:spPr'))
    for old in spPr.findall(qn('a:solidFill')): spPr.remove(old)
    for old in spPr.findall(qn('a:gradFill')):  spPr.remove(old)
    h1 = f"{c1[0]:02x}{c1[1]:02x}{c1[2]:02x}"
    h2 = f"{c2[0]:02x}{c2[1]:02x}{c2[2]:02x}"
    xml = f"""<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="{h1}"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="{h2}"/></a:gs>
  </a:gsLst>
  <a:lin ang="{angle*60000}" scaled="0"/>
</a:gradFill>"""
    spPr.insert(len(list(spPr)), parse_xml(xml))
    return shape

def add_overlay(slide, l, t, w, h, color, opacity=0.80):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    sp   = shape._element
    spPr = sp.find(qn('p:spPr'))
    for old in spPr.findall(qn('a:solidFill')): spPr.remove(old)
    for old in spPr.findall(qn('a:gradFill')):  spPr.remove(old)
    ch = f"{color[0]:02x}{color[1]:02x}{color[2]:02x}"
    av = int(opacity * 100000)
    xml = f"""<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:srgbClr val="{ch}"><a:alpha val="{av}"/></a:srgbClr>
</a:solidFill>"""
    spPr.insert(len(list(spPr)), parse_xml(xml))
    return shape

# ─── Animation helpers ────────────────────────────────────────────────────────
_anim_id_counter = [10]

def _next_id():
    _anim_id_counter[0] += 1
    return _anim_id_counter[0]

def _get_childTnLst2(slide):
    spTree = slide.shapes._spTree
    sld    = spTree.getparent()
    timing = sld.find(qn('p:timing'))
    if timing is None:
        timing = etree.SubElement(sld, qn('p:timing'))
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = etree.SubElement(timing, qn('p:tnLst'))
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = etree.SubElement(tnLst, qn('p:par'))
    cTn = par.find(qn('p:cTn'))
    if cTn is None:
        cTn = etree.SubElement(par, qn('p:cTn'))
        cTn.set('id', '1'); cTn.set('dur', 'indefinite')
        cTn.set('restart', 'whenNotActive'); cTn.set('nodeType', 'tmRoot')
    cTnLst = cTn.find(qn('p:childTnLst'))
    if cTnLst is None:
        cTnLst = etree.SubElement(cTn, qn('p:childTnLst'))
    seq = cTnLst.find(qn('p:seq'))
    if seq is None:
        seq = etree.SubElement(cTnLst, qn('p:seq'))
        seq.set('concurrent', '1'); seq.set('nextAc', 'seek')
    cTn2 = seq.find(qn('p:cTn'))
    if cTn2 is None:
        cTn2 = etree.SubElement(seq, qn('p:cTn'))
        cTn2.set('id', '2'); cTn2.set('dur', 'indefinite'); cTn2.set('nodeType', 'mainSeq')
    cTnLst2 = cTn2.find(qn('p:childTnLst'))
    if cTnLst2 is None:
        cTnLst2 = etree.SubElement(cTn2, qn('p:childTnLst'))
    if seq.find(qn('p:prevCondLst')) is None:
        pcl = etree.SubElement(seq, qn('p:prevCondLst'))
        c = etree.SubElement(pcl, qn('p:cond')); c.set('evt','onPrevClick'); c.set('delay','0')
        t = etree.SubElement(c, qn('p:tgtEl')); etree.SubElement(t, qn('p:sldTgt'))
    if seq.find(qn('p:nextCondLst')) is None:
        ncl = etree.SubElement(seq, qn('p:nextCondLst'))
        c = etree.SubElement(ncl, qn('p:cond')); c.set('evt','onNextClick'); c.set('delay','0')
        t = etree.SubElement(c, qn('p:tgtEl')); etree.SubElement(t, qn('p:sldTgt'))
    return cTnLst2

def add_fly_in(slide, shape, delay_ms=0):
    sp_id = shape._element.get('id') or shape._element.find('.//' + qn('p:sp'))
    if hasattr(shape, '_element'):
        el = shape._element
        nvSpPr = el.find('.//' + qn('p:nvSpPr')) or el.find('.//' + qn('p:nvPicPr')) or el.find('.//' + qn('p:nvGrpSpPr'))
        cNvPr  = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
        sid    = cNvPr.get('id') if cNvPr is not None else '1'
    else:
        sid = '1'
    cTnLst2 = _get_childTnLst2(slide)
    pid = str(_next_id()); cid = str(_next_id()); eid = str(_next_id())
    xml = f"""<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cTn id="{pid}" fill="hold">
    <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{cid}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:animEffect xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          transition="in" filter="fly">
              <p:cBhvr>
                <p:cTn id="{eid}" dur="600" fill="hold"/>
                <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
              </p:cBhvr>
            </p:animEffect>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>"""
    cTnLst2.append(parse_xml(xml))

def add_fade(slide, shape, delay_ms=0):
    if hasattr(shape, '_element'):
        el = shape._element
        nvSpPr = el.find('.//' + qn('p:nvSpPr')) or el.find('.//' + qn('p:nvPicPr'))
        cNvPr  = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
        sid    = cNvPr.get('id') if cNvPr is not None else '1'
    else:
        sid = '1'
    cTnLst2 = _get_childTnLst2(slide)
    pid = str(_next_id()); cid = str(_next_id()); eid = str(_next_id())
    xml = f"""<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cTn id="{pid}" fill="hold">
    <p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{cid}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:animEffect xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          transition="in" filter="fade">
              <p:cBhvr>
                <p:cTn id="{eid}" dur="500" fill="hold"/>
                <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
              </p:cBhvr>
            </p:animEffect>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>"""
    cTnLst2.append(parse_xml(xml))

# ─── Decorative helpers ───────────────────────────────────────────────────────
def top_accent_bar(s, color=GREEN):
    add_rect(s, 0, 0, 13.33, 0.07, fill_color=color)

def bottom_footer(s, text="Why BPL Fails to Attract Urban Youth  |  Batch 232 — Textile Department"):
    add_rect(s, 0, 7.25, 13.33, 0.25, fill_color=RGBColor(0x08, 0x10, 0x22))
    add_text(s, f"© Copyright by Tahorim Somrat  •  {text}", 0.2, 7.26, 12.9, 0.24,
             font_size=8.5, color=MID_GRAY, align=PP_ALIGN.CENTER)

def slide_badge(s, label, l=0.3, t=0.12, w=3.0, color=ACCENT_BLUE):
    add_rect(s, l, t, w, 0.42, fill_color=color)
    add_text(s, label, l, t, w, 0.42, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def slide_number(s, num):
    add_text(s, f"{num:02d}", 12.6, 0.12, 0.6, 0.4,
             font_size=20, bold=True, color=RGBColor(0x2A, 0x3A, 0x5A))

def divider(s, t=3.0, color=RGBColor(0x22, 0x33, 0x55)):
    add_rect(s, 0.5, t, 12.33, 0.04, fill_color=color)

def section_title(s, text, t=0.65, color=WHITE, size=30):
    tx = add_text(s, text, 0.5, t, 12.33, 0.85, font_size=size, bold=True,
                  color=color, align=PP_ALIGN.CENTER)
    add_fly_in(s, tx, 100)

def underline_bar(s, color=GREEN, t=1.5, w=3.0, l=5.16):
    add_rect(s, l, t, w, 0.06, fill_color=color)

def add_icon_text_row(s, icon, text, l, t, icon_color=GREEN, text_color=WHITE, size=16):
    add_text(s, icon, l, t, 0.5, 0.45, font_size=size+2, bold=True, color=icon_color)
    add_text(s, text, l+0.5, t, 10.0, 0.45, font_size=size, color=text_color)

def stat_card(s, l, t, w, h, number, label, unit="", num_color=GOLD, bg_color=DARK_CARD, border_color=GOLD):
    add_rect(s, l, t, w, h, fill_color=bg_color, line_color=border_color, line_width=1.5)
    add_text(s, number+unit, l, t+0.1, w, h*0.55, font_size=34, bold=True,
             color=num_color, align=PP_ALIGN.CENTER)
    add_text(s, label, l, t+h*0.55, w, h*0.4, font_size=11, bold=False,
             color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ─── Slide Functions ──────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Slide 1 — Title Slide"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)

    # Background image
    if os.path.exists(IMG_SLIDE1):
        s.shapes.add_picture(IMG_SLIDE1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    # Dark overlay for readability
    add_overlay(s, 0, 0, 13.33, 7.5, DARK_BG, opacity=0.72)

    # Top decorative gradient bar
    add_gradient_rect(s, 0, 0, 13.33, 0.12, GREEN, ACCENT_BLUE, angle=0)

    # Title block background card
    add_overlay(s, 1.2, 0.7, 10.93, 3.5, DARK_CARD, opacity=0.88)
    add_rect(s, 1.2, 0.7, 10.93, 3.5, fill_color=None, line_color=GREEN, line_width=1.5)

    # Main title
    t1 = add_text(s, "Why Bangladesh Premier League", 1.5, 0.9, 10.33, 0.95,
                  font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_fly_in(s, t1, 0)

    t2 = add_text(s, "Fails to Attract Urban Youth", 1.5, 1.75, 10.33, 0.85,
                  font_size=34, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_fly_in(s, t2, 200)

    # Subtitle underline
    add_gradient_rect(s, 3.5, 2.62, 6.33, 0.06, GREEN, ACCENT_BLUE, angle=0)

    t3 = add_text(s, "A Research Presentation on BPL Fan Engagement & Youth Disconnect", 1.5, 2.75, 10.33, 0.55,
                  font_size=15, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_fly_in(s, t3, 350)

    # Department / Batch badge
    dept_bg = add_rect(s, 4.16, 3.45, 5.0, 0.6, fill_color=RGBColor(0x00, 0x55, 0x1A))
    add_text(s, "Batch 232  —  Textile Department", 4.16, 3.45, 5.0, 0.6,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Presented by card
    add_overlay(s, 0.5, 4.35, 5.5, 2.7, DARK_CARD, opacity=0.90)
    add_rect(s, 0.5, 4.35, 5.5, 2.7, fill_color=None, line_color=ACCENT_BLUE, line_width=1.2)
    add_text(s, "PRESENTED BY", 0.7, 4.45, 5.1, 0.45,
             font_size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    divider(s, t=4.9, color=RGBColor(0x1A, 0x2A, 0x4A))

    members = [
        ("2230800230", "Mahmudul Hasan"),
        ("2230800231", "Md. Rahat"),
        ("2230800232", "Md. Rifat"),
        ("2230800233", "Md. Tanvir"),
        ("2230800234", "Md. Towhid"),
    ]
    for i, (id_, name) in enumerate(members):
        add_text(s, f"◆  {id_}  —  {name}", 0.7, 5.0 + i*0.37, 5.1, 0.38,
                 font_size=12.5, color=WHITE)

    # Supervisor side note
    add_overlay(s, 7.3, 4.35, 5.5, 2.7, DARK_CARD, opacity=0.88)
    add_rect(s, 7.3, 4.35, 5.5, 2.7, fill_color=None, line_color=GOLD, line_width=1.2)
    add_text(s, "TOPIC OVERVIEW", 7.5, 4.45, 5.1, 0.45,
             font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    divider(s, t=4.9, color=RGBColor(0x33, 0x28, 0x00))
    overview = [
        "⚽  BPL Fan Engagement",
        "📊  Statistical Comparison",
        "🎯  Root Causes & Analysis",
        "💡  Actionable Recommendations",
        "🏆  Vision for BPL's Future",
    ]
    for i, ov in enumerate(overview):
        add_text(s, ov, 7.5, 5.0 + i*0.37, 5.0, 0.38, font_size=12.5, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_02_reality(prs):
    """Slide 2 — The Reality"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s)
    slide_badge(s, "THE REALITY", color=RED)
    slide_number(s, 2)

    t1 = add_text(s, "Bangladesh Loves Football ❤️", 0.5, 0.65, 12.33, 0.75,
                  font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_fly_in(s, t1, 0)

    but_bg = add_rect(s, 4.66, 1.42, 4.0, 0.6, fill_color=RGBColor(0x3A, 0x10, 0x00))
    t2 = add_text(s, "B U T . . .", 4.66, 1.42, 4.0, 0.6,
                  font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_fly_in(s, t2, 300)

    t3 = add_text(s, "Bangladesh Doesn't Watch Local Football ❌", 0.5, 2.1, 12.33, 0.7,
                  font_size=26, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_fly_in(s, t3, 550)

    divider(s, t=2.95)

    # ── European card ──
    add_gradient_rect(s, 0.4, 3.1, 5.9, 3.8, RGBColor(0x03,0x1A,0x0A), DARK_CARD, angle=135)
    add_rect(s, 0.4, 3.1, 5.9, 3.8, fill_color=None, line_color=GREEN, line_width=2.0)
    add_text(s, "🌍  European Football", 0.6, 3.18, 5.5, 0.55,
             font_size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    divider(s, t=3.73, color=RGBColor(0x08,0x33,0x15))
    euro_pts = ["✔  Millions of Bangladeshi Fans watching",
                "✔  Enormous Social Media Buzz daily",
                "✔  Packed stadiums & electric atmosphere",
                "✔  Global superstars & brand power",
                "✔  High-quality TV production & coverage"]
    for i, pt in enumerate(euro_pts):
        add_text(s, pt, 0.6, 3.82 + i*0.55, 5.5, 0.52, font_size=14, color=LIGHT_GRAY)

    # ── BPL card ──
    add_gradient_rect(s, 7.0, 3.1, 5.9, 3.8, RGBColor(0x1A,0x03,0x03), DARK_CARD, angle=135)
    add_rect(s, 7.0, 3.1, 5.9, 3.8, fill_color=None, line_color=RED, line_width=2.0)
    add_text(s, "🇧🇩  Bangladesh Premier League", 7.2, 3.18, 5.5, 0.55,
             font_size=17, bold=True, color=RED, align=PP_ALIGN.CENTER)
    divider(s, t=3.73, color=RGBColor(0x33,0x08,0x08))
    bpl_pts = ["❌  Very low stadium attendance",
               "❌  Minimal social media presence",
               "❌  Limited broadcast coverage",
               "❌  No internationally known star players",
               "❌  Weak marketing & fan engagement"]
    for i, pt in enumerate(bpl_pts):
        add_text(s, pt, 7.2, 3.82 + i*0.55, 5.5, 0.52, font_size=14, color=LIGHT_GRAY)

    # VS badge
    add_gradient_rect(s, 6.16, 4.3, 1.0, 0.9, RED, GREEN, angle=90)
    add_text(s, "VS", 6.16, 4.3, 1.0, 0.9, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bottom_footer(s)


def slide_03_stats(prs):
    """Slide 3 — Key Statistics"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, GOLD)
    slide_badge(s, "KEY STATISTICS", color=DARK_GREEN)
    slide_number(s, 3)

    section_title(s, "The Numbers Don't Lie 📊", size=28)
    underline_bar(s, GOLD, t=1.5, w=4.0, l=4.66)

    add_text(s, "Bangladesh vs. European Football — Viewership & Engagement Data",
             0.5, 1.62, 12.33, 0.5, font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Stat cards row 1
    stats1 = [
        ("90%",  "Bangladeshis follow at\nleast one European club", GOLD),
        ("12%",  "Follow BPL clubs\nactively on social media", RED),
        ("3M+",  "Bangladeshis watch EPL\nlive every week", GREEN),
    ]
    for i, (num, lbl, col) in enumerate(stats1):
        stat_card(s, 0.5 + i*4.3, 2.2, 3.9, 2.0, num, lbl, num_color=col,
                  bg_color=DARK_CARD, border_color=col)

    # Stat cards row 2
    stats2 = [
        ("8,000",  "Average BPL stadium\nattendance per match", ORANGE),
        ("60,000", "Stadium capacity rarely\nused — only 13% filled", RED),
        ("< 5%",   "Youth aged 18–30 who\nwatch BPL regularly", TEAL),
    ]
    for i, (num, lbl, col) in enumerate(stats2):
        stat_card(s, 0.5 + i*4.3, 4.35, 3.9, 2.0, num, lbl, num_color=col,
                  bg_color=DARK_CARD, border_color=col)

    add_text(s, "* Sources: Bangladesh Football Federation (BFF), Statista, Sports Research BD",
             0.5, 6.8, 12.33, 0.35, font_size=8.5, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)
    bottom_footer(s)


def slide_04_why_euro(prs):
    """Slide 4 — Why Youth Prefer European Football"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, ACCENT_BLUE)
    slide_badge(s, "WHY EUROPEAN FOOTBALL?", color=ACCENT_BLUE)
    slide_number(s, 4)

    section_title(s, "5 Reasons Youth Prefer European Football", size=27)
    underline_bar(s, ACCENT_BLUE, t=1.5, w=4.5, l=4.41)

    # Optional image
    if os.path.exists(IMG_INTL):
        try:
            s.shapes.add_picture(IMG_INTL, Inches(8.5), Inches(1.7), Inches(4.5), Inches(5.3))
            add_overlay(s, 8.5, 1.7, 4.5, 5.3, DARK_BG, opacity=0.18)
            add_rect(s, 8.5, 1.7, 4.5, 5.3, fill_color=None, line_color=ACCENT_BLUE, line_width=1.5)
        except Exception:
            pass

    reasons = [
        ("1", "⚡  Superstar Power", "Ronaldo, Messi, Mbappé — icons with 500M+ social followers", GOLD),
        ("2", "🎬  World-Class Broadcast", "4K coverage, expert commentary, UEFA app & highlights", ACCENT_BLUE),
        ("3", "🏟️  Electrifying Atmosphere", "50,000+ packed stadiums, ultras, flares, chants worldwide", GREEN),
        ("4", "📱  Digital & Social Media", "Viral clips, TikTok, YouTube — content available 24/7", TEAL),
        ("5", "🌐  Global Brand Identity", "Club shirts, merchandise, FIFA games — culture & lifestyle", ORANGE),
    ]
    for i, (num, title, desc, col) in enumerate(reasons):
        ty = 1.75 + i * 1.0
        add_gradient_rect(s, 0.4, ty, 7.8, 0.85, DARK_CARD, RGBColor(0x0C,0x14,0x28), angle=0)
        add_rect(s, 0.4, ty, 0.5, 0.85, fill_color=col)
        add_text(s, num, 0.4, ty, 0.5, 0.85, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.05, ty + 0.05, 3.5, 0.4, font_size=15, bold=True, color=col)
        add_text(s, desc,  1.05, ty + 0.42, 6.9, 0.4, font_size=12.5, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_05_bpl_struggles(prs):
    """Slide 5 — Why BPL Struggles"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, RED)
    slide_badge(s, "BPL STRUGGLES", color=RED)
    slide_number(s, 5)

    section_title(s, "Why Bangladesh Premier League Struggles ⚠️", size=26, color=RED)
    underline_bar(s, RED, t=1.5, w=4.0, l=4.66)

    issues = [
        ("📣", "Poor Marketing",        "No star promotions, no hype campaigns, no youth outreach programs"),
        ("📺", "Weak Broadcasting",     "Inconsistent TV schedule, poor stream quality, no streaming app"),
        ("🌟", "No Star Players",       "Local players lack global recognition; no foreign marquee signings"),
        ("🏟️", "Stadium Experience",    "Outdated venues, poor facilities, security concerns deter attendance"),
        ("💰", "Financial Instability", "Clubs face funding issues; salaries delayed; no investor confidence"),
        ("📱", "Digital Absence",       "Almost zero social media strategy; no highlight reels or shorts"),
    ]

    for i, (icon, title, desc) in enumerate(issues):
        col = i % 2
        row = i // 2
        lx = 0.4  + col * 6.5
        ty = 1.75 + row * 1.75

        add_gradient_rect(s, lx, ty, 6.1, 1.55, RGBColor(0x1A,0x05,0x05), DARK_CARD, angle=135)
        add_rect(s, lx, ty, 6.1, 1.55, fill_color=None, line_color=RED, line_width=1.3)
        add_text(s, icon, lx+0.1, ty+0.1, 0.65, 0.65, font_size=24)
        add_text(s, title, lx+0.85, ty+0.12, 4.9, 0.45, font_size=15, bold=True, color=RED)
        add_text(s, desc,  lx+0.85, ty+0.6,  4.9, 0.85, font_size=11.5, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_06_bd_players(prs):
    """Slide 6 — Bangladesh Players"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, GREEN)
    slide_badge(s, "BANGLADESHI FOOTBALL STARS", color=DARK_GREEN)
    slide_number(s, 6)

    section_title(s, "Bangladesh Has Talented Players — But Lacks Exposure 🇧🇩", size=24, color=GREEN)
    underline_bar(s, GREEN, t=1.5, w=5.0, l=4.16)

    # Player image if available
    if os.path.exists(IMG_BD):
        try:
            s.shapes.add_picture(IMG_BD, Inches(9.0), Inches(1.65), Inches(4.0), Inches(5.4))
            add_overlay(s, 9.0, 1.65, 4.0, 5.4, DARK_BG, opacity=0.15)
            add_rect(s, 9.0, 1.65, 4.0, 5.4, fill_color=None, line_color=GREEN, line_width=1.5)
        except Exception:
            pass

    players = [
        ("🟢", "Hamza Choudhury", "British-Bangladeshi midfielder, Leicester City FC\nFirst Bangladeshi-origin player in Premier League"),
        ("🟢", "Jamal Bhuyan",    "Bangladesh national team captain & midfielder\nPlayed in professional leagues in Denmark & Spain"),
        ("🟢", "Sheikh Morsalin", "Top scorer of Bangladesh national team\nKnown for explosive speed and clinical finishing"),
        ("🟢", "Biplu Ahmed",     "Experienced winger for national team\nBFF Player of the Year multiple times"),
        ("🟢", "Rakib Hossain",   "Young talent from Abahani FC\nConsidered the next big prospect of Bangladeshi football"),
    ]

    for i, (icon, name, desc) in enumerate(players):
        ty = 1.75 + i * 1.03
        add_gradient_rect(s, 0.4, ty, 8.2, 0.9, DARK_CARD, RGBColor(0x03,0x1A,0x0A), angle=0)
        add_rect(s, 0.4, ty, 0.08, 0.9, fill_color=GREEN)
        add_text(s, icon, 0.6, ty+0.05, 0.5, 0.5, font_size=18)
        add_text(s, name, 1.2, ty+0.05, 4.5, 0.4, font_size=15, bold=True, color=GREEN)
        add_text(s, desc, 1.2, ty+0.45, 6.8, 0.45, font_size=11, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_07_intl_stars(prs):
    """Slide 7 — International Stars & Comparison"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, GOLD)
    slide_badge(s, "INTERNATIONAL STARS", color=RGBColor(0xAA,0x77,0x00))
    slide_number(s, 7)

    section_title(s, "Why Global Stars Drive Fan Engagement 🌟", size=26, color=GOLD)
    underline_bar(s, GOLD, t=1.5, w=4.5, l=4.41)

    if os.path.exists(IMG_INTL):
        try:
            s.shapes.add_picture(IMG_INTL, Inches(0.4), Inches(1.7), Inches(4.3), Inches(5.3))
            add_overlay(s, 0.4, 1.7, 4.3, 5.3, DARK_BG, opacity=0.12)
            add_rect(s, 0.4, 1.7, 4.3, 5.3, fill_color=None, line_color=GOLD, line_width=1.5)
        except Exception:
            pass

    stars = [
        ("⭐ Cristiano Ronaldo", "635M+ Instagram followers\nGlobal brand worth $1.5 billion", GOLD),
        ("⭐ Lionel Messi",      "500M+ followers; Inter Miami\nArgentina World Cup hero 2022", ORANGE),
        ("⭐ Kylian Mbappé",     "Real Madrid star; fastest in EPL era\nFrance's golden generation leader", ACCENT_BLUE),
        ("⭐ Erling Haaland",    "52 goals/season — Man City striker\nMost watched striker in the world", GREEN),
    ]

    for i, (name, desc, col) in enumerate(stars):
        ty = 1.75 + i * 1.35
        add_gradient_rect(s, 5.0, ty, 7.9, 1.2, DARK_CARD, RGBColor(0x12,0x10,0x05), angle=0)
        add_rect(s, 5.0, ty, 7.9, 1.2, fill_color=None, line_color=col, line_width=1.3)
        add_text(s, name, 5.2, ty+0.08, 7.4, 0.45, font_size=15, bold=True, color=col)
        add_text(s, desc, 5.2, ty+0.55, 7.4, 0.58, font_size=12, color=LIGHT_GRAY)

    add_text(s, "💡 BPL needs a similar star culture to attract youth viewers!", 0.5, 7.0, 12.33, 0.38,
             font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    bottom_footer(s)


def slide_08_marketing(prs):
    """Slide 8 — Marketing Failure"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, ORANGE)
    slide_badge(s, "MARKETING FAILURE", color=ORANGE)
    slide_number(s, 8)

    section_title(s, "BPL's Marketing vs. European Football Marketing 📣", size=25, color=ORANGE)
    underline_bar(s, ORANGE, t=1.5, w=5.0, l=4.16)

    # Two column comparison
    headers = [("BPL Marketing ❌", RED, 0.4, 7.1),
               ("European Football ✔", GREEN, 7.1, 5.8)]
    bpl_pts = [
        "No official promotion on social media",
        "Zero youth ambassador programs",
        "No jersey culture or merchandise stores",
        "Matches not promoted on mainstream TV",
        "No fan zones or pre-match events",
    ]
    eur_pts = [
        "100M+ ad spend per club per season",
        "Global ambassador partnerships",
        "Official jerseys sold in 180+ countries",
        "Prime time slots + dedicated channels",
        "Fan experience zones in every stadium",
    ]

    for j, (label, col, lx, lw) in enumerate(headers):
        add_gradient_rect(s, lx, 2.0, lw, 4.8, DARK_CARD, DARK_BG, angle=90)
        add_rect(s, lx, 2.0, lw, 4.8, fill_color=None, line_color=col, line_width=1.8)
        add_text(s, label, lx+0.1, 2.05, lw-0.2, 0.5, font_size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
        pts = bpl_pts if j == 0 else eur_pts
        for i, pt in enumerate(pts):
            pre = "❌ " if j == 0 else "✅ "
            add_text(s, pre + pt, lx+0.2, 2.65 + i*0.78, lw-0.3, 0.65, font_size=12.5, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_09_broadcasting(prs):
    """Slide 9 — Broadcasting Problem"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, TEAL)
    slide_badge(s, "BROADCASTING FAILURE", color=RGBColor(0x00,0x88,0x77))
    slide_number(s, 9)

    section_title(s, "The Broadcasting Gap 📺", size=30, color=TEAL)
    underline_bar(s, TEAL, t=1.5, w=3.5, l=4.91)

    points = [
        ("📺", "No Streaming Platform", "BPL has no official app or OTT platform; European leagues have Prime Video, DAZN, Peacock, Disney+"),
        ("🎙️", "Poor Commentary",       "Local broadcasting lacks professional commentary; no English or multilingual options"),
        ("📷", "Camera Quality",         "Single-camera coverage vs. UEFA's 30+ camera, VR, 4K drone coverage"),
        ("⏰", "Irregular Scheduling",   "Match times change last minute; no fixed schedule disrupts fan planning"),
        ("🌍", "No International Feed",  "BPL matches unavailable outside Bangladesh; zero global audience growth"),
    ]

    for i, (icon, title, desc) in enumerate(points):
        ty = 1.75 + i * 1.05
        add_gradient_rect(s, 0.4, ty, 12.5, 0.9, DARK_CARD, RGBColor(0x03,0x16,0x14), angle=0)
        add_rect(s, 0.4, ty, 0.7, 0.9, fill_color=TEAL)
        add_text(s, icon, 0.4, ty+0.12, 0.7, 0.65, font_size=20, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.25, ty+0.06, 3.0, 0.4, font_size=15, bold=True, color=TEAL)
        add_text(s, desc, 1.25, ty+0.48, 11.3, 0.42, font_size=11.5, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_10_stadium(prs):
    """Slide 10 — Stadium Experience"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, ORANGE)
    slide_badge(s, "STADIUM EXPERIENCE", color=ORANGE)
    slide_number(s, 10)

    if os.path.exists(IMG_STADIUM):
        try:
            s.shapes.add_picture(IMG_STADIUM, Inches(7.5), Inches(0.1), Inches(5.83), Inches(7.4))
            add_overlay(s, 7.5, 0.1, 5.83, 7.4, DARK_BG, opacity=0.55)
        except Exception:
            pass

    section_title(s, "Stadium Experience: Night & Day Difference 🏟️", size=24, color=ORANGE)
    underline_bar(s, ORANGE, t=1.5, w=5.5, l=0.5)

    issues = [
        ("🚧", "Poor Infrastructure",   "Old, poorly maintained stands with broken seats"),
        ("🔒", "Security Concerns",     "Crowd management issues scare away families & women"),
        ("🍔", "No Food/Entertainment", "No food courts, fan zones, or entertainment areas"),
        ("🚌", "Transport Access",      "Lack of easy transport to stadiums in major cities"),
        ("🌧️", "Weather Exposure",      "No covered seating areas; fans leave in rain"),
        ("💡", "Poor Lighting",         "Outdated lighting fails night broadcast standards"),
    ]

    for i, (icon, title, desc) in enumerate(issues):
        ty = 1.75 + i * 0.92
        add_gradient_rect(s, 0.4, ty, 6.8, 0.8, DARK_CARD, DARK_BG, angle=0)
        add_rect(s, 0.4, ty, 0.6, 0.8, fill_color=ORANGE)
        add_text(s, icon, 0.4, ty+0.07, 0.6, 0.6, font_size=18, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.12, ty+0.06, 2.5, 0.35, font_size=13, bold=True, color=ORANGE)
        add_text(s, desc,  1.12, ty+0.42, 5.7, 0.35, font_size=11, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_11_youth_disconnect(prs):
    """Slide 11 — Youth Disconnect"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, ACCENT_BLUE)
    slide_badge(s, "YOUTH DISCONNECT", color=ACCENT_BLUE)
    slide_number(s, 11)

    section_title(s, "Why Urban Youth Are Disconnected from BPL 📱", size=25, color=ACCENT_BLUE)
    underline_bar(s, ACCENT_BLUE, t=1.5, w=5.0, l=4.16)

    # Quote block
    add_gradient_rect(s, 0.5, 1.75, 12.33, 1.1, RGBColor(0x08,0x15,0x2A), DARK_CARD, angle=0)
    add_rect(s, 0.5, 1.75, 0.08, 1.1, fill_color=ACCENT_BLUE)
    add_text(s, '"I watch UCL highlights the moment they drop — BPL? I don\'t even know when the season starts."',
             0.75, 1.82, 11.8, 0.95, font_size=14, italic=True, color=LIGHT_GRAY)
    add_text(s, "— Urban Bangladeshi Youth (Survey Response, 2024)", 0.75, 2.72, 11.8, 0.28,
             font_size=10.5, color=MID_GRAY, italic=True)

    reasons = [
        ("📱", "No Social Media Content",    "BPL clubs barely post; no reels, highlights, or memes"),
        ("🎮", "Gaming & FIFA Culture",       "Youth identify with EPL/LaLiga clubs through FIFA video games"),
        ("🧢", "No Fan Identity",             "BPL clubs have no cool jersey, badge, or fan culture youth relate to"),
        ("📺", "Content Unavailability",      "No YouTube channel, Shorts, or TikTok presence for BPL"),
        ("🤷", "No Role Model Players",       "Youth can't name a BPL star the way they name Messi or Ronaldo"),
    ]

    for i, (icon, title, desc) in enumerate(reasons):
        ty = 3.1 + i * 0.83
        add_gradient_rect(s, 0.4, ty, 12.5, 0.72, DARK_CARD, RGBColor(0x06,0x10,0x22), angle=0)
        add_rect(s, 0.4, ty, 0.55, 0.72, fill_color=ACCENT_BLUE)
        add_text(s, icon, 0.4, ty+0.08, 0.55, 0.55, font_size=18, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.1, ty+0.05, 3.5, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
        add_text(s, desc, 1.1, ty+0.4, 11.4, 0.32, font_size=11.5, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_12_recommendations(prs):
    """Slide 12 — Recommendations"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    top_accent_bar(s, TEAL)
    slide_badge(s, "RECOMMENDATIONS", color=RGBColor(0x00,0x77,0x66))
    slide_number(s, 12)

    section_title(s, "How to Fix BPL & Attract Urban Youth 💡", size=26, color=TEAL)
    underline_bar(s, TEAL, t=1.5, w=4.5, l=4.41)

    recs = [
        ("1", "Digital Revolution",     "Launch official OTT app, YouTube channel, TikTok, Instagram reels for highlights", TEAL),
        ("2", "Star Player Signings",   "Bring in 2-3 international marquee players to spark excitement and media attention", GOLD),
        ("3", "Youth Engagement",       "Free student tickets, campus fan clubs, BPL gaming tournaments & fantasy leagues", GREEN),
        ("4", "Upgrade Stadiums",       "Modern covered stands, fan zones, food courts, LED screens, proper security", ORANGE),
        ("5", "Marketing Campaigns",    "Billboard, TV, social media campaigns targeting 18–30 year old urban demographic", ACCENT_BLUE),
        ("6", "Professional Broadcasting", "Partner with international broadcasters; 4K coverage, multi-language commentary", RED),
    ]

    for i, (num, title, desc, col) in enumerate(recs):
        col_idx = i % 2
        row_idx = i // 2
        lx = 0.4 + col_idx * 6.5
        ty = 1.75 + row_idx * 1.75

        add_gradient_rect(s, lx, ty, 6.1, 1.55, DARK_CARD, DARK_BG, angle=90)
        add_rect(s, lx, ty, 6.1, 1.55, fill_color=None, line_color=col, line_width=1.5)
        add_rect(s, lx, ty, 0.55, 1.55, fill_color=col)
        add_text(s, num, lx, ty, 0.55, 1.55, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, lx+0.65, ty+0.1, 5.1, 0.45, font_size=14, bold=True, color=col)
        add_text(s, desc,  lx+0.65, ty+0.6, 5.1, 0.85, font_size=11, color=LIGHT_GRAY)

    bottom_footer(s)


def slide_13_vision(prs):
    """Slide 13 — Vision for BPL"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)

    # Dramatic gradient bg
    add_gradient_rect(s, 0, 0, 13.33, 7.5, RGBColor(0x03,0x1A,0x0A), DARK_BG, angle=90)
    top_accent_bar(s, GREEN)
    slide_badge(s, "THE VISION", color=DARK_GREEN)
    slide_number(s, 13)

    section_title(s, "A Future Where BPL Inspires a Nation 🏆", size=28, color=GREEN)
    underline_bar(s, GOLD, t=1.5, w=4.5, l=4.41)

    vision_points = [
        "🏟️  30,000+ average stadium attendance per match",
        "📱  1 Million+ BPL social media followers by 2026",
        "🌍  International broadcast deals with Al Jazeera, Willow TV",
        "⚽  Bangladeshi players competing in Asian Champions League",
        "👕  BPL jersey culture among youth across Dhaka, Chittagong",
        "💰  Major brand sponsorships: Grameenphone, Bashundhara, bKash",
    ]

    add_overlay(s, 0.5, 1.75, 12.33, 5.1, DARK_CARD, opacity=0.70)
    add_rect(s, 0.5, 1.75, 12.33, 5.1, fill_color=None, line_color=GREEN, line_width=1.5)

    for i, pt in enumerate(vision_points):
        ty = 2.0 + i * 0.78
        add_text(s, pt, 1.0, ty, 11.2, 0.65, font_size=15.5, color=LIGHT_GRAY)
        add_rect(s, 0.65, ty + 0.2, 0.15, 0.25, fill_color=GREEN)

    add_text(s, "\"Football is not just a sport — it is the heartbeat of a nation.\"",
             0.5, 6.9, 12.33, 0.38, font_size=12.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    bottom_footer(s)


def slide_14_conclusion(prs):
    """Slide 14 — Conclusion"""
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)

    # Bg overlay
    if os.path.exists(IMG_SLIDE1):
        try:
            s.shapes.add_picture(IMG_SLIDE1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        except Exception:
            pass
    add_overlay(s, 0, 0, 13.33, 7.5, DARK_BG, opacity=0.80)

    add_gradient_rect(s, 0, 0, 13.33, 0.12, RED, GREEN, angle=0)
    slide_badge(s, "CONCLUSION", color=DARK_GREEN)
    slide_number(s, 14)

    section_title(s, "Conclusion", size=32, color=WHITE)
    underline_bar(s, GREEN, t=1.52, w=2.5, l=5.41)

    add_overlay(s, 1.0, 1.75, 11.33, 3.6, DARK_CARD, opacity=0.88)
    add_rect(s, 1.0, 1.75, 11.33, 3.6, fill_color=None, line_color=GREEN, line_width=2.0)

    add_text(s, "Bangladesh Premier League's failure to attract urban youth", 1.3, 1.9, 10.73, 0.65,
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "is not a football problem — it is a MANAGEMENT & MARKETING problem.", 1.3, 2.5, 10.73, 0.65,
             font_size=17, bold=True, color=RED, align=PP_ALIGN.CENTER)

    divider(s, t=3.2, color=RGBColor(0x22,0x33,0x22))

    add_text(s, "✅  With proper digital marketing, quality broadcasting,\n"
                "      modern stadiums, and star-player culture —\n"
                "      BPL CAN become the pride of Bangladeshi youth.",
             1.5, 3.3, 10.0, 1.1, font_size=15.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Final tagline
    add_gradient_rect(s, 1.5, 5.55, 10.33, 0.75, DARK_GREEN, RGBColor(0x00,0x55,0x1A), angle=0)
    add_text(s, "🇧🇩  BPL — From Local League to National Pride  🇧🇩", 1.5, 5.55, 10.33, 0.75,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "Thank You for Your Attention!", 0.5, 6.4, 12.33, 0.5,
             font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    bottom_footer(s)


# ─── Main ────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    print("[1/14] Slide 1 — Title...")
    slide_01_title(prs)
    print("[2/14] Slide 2 — The Reality...")
    slide_02_reality(prs)
    print("[3/14] Slide 3 — Statistics...")
    slide_03_stats(prs)
    print("[4/14] Slide 4 — Why Euro Football...")
    slide_04_why_euro(prs)
    print("[5/14] Slide 5 — BPL Struggles...")
    slide_05_bpl_struggles(prs)
    print("[6/14] Slide 6 — BD Players...")
    slide_06_bd_players(prs)
    print("[7/14] Slide 7 — International Stars...")
    slide_07_intl_stars(prs)
    print("[8/14] Slide 8 — Marketing...")
    slide_08_marketing(prs)
    print("[9/14] Slide 9 — Broadcasting...")
    slide_09_broadcasting(prs)
    print("[10/14] Slide 10 — Stadium Experience...")
    slide_10_stadium(prs)
    print("[11/14] Slide 11 — Youth Disconnect...")
    slide_11_youth_disconnect(prs)
    print("[12/14] Slide 12 — Recommendations...")
    slide_12_recommendations(prs)
    print("[13/14] Slide 13 — Vision...")
    slide_13_vision(prs)
    print("[14/14] Slide 14 — Conclusion...")
    slide_14_conclusion(prs)

    prs.save(OUT_PATH)
    print(f"\n[DONE] Presentation saved: {OUT_PATH}")
    print(f"       Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
